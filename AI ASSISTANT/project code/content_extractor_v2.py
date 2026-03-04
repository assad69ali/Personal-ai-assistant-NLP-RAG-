
import json
import fitz  # PyMuPDF
import nbformat
import pandas as pd
from pathlib import Path
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup

# === Paths ===
BASE_DIR = Path(r"D:\project ir\dataset for ir")
METADATA_PATH = Path(r"D:\project ir\indexed_metadata.json")
OUTPUT_PATH = Path(r"D:\project ir\enriched_with_text.json")

def read_text_file(path):
    try:
        return Path(path).read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        return f"[ERROR reading text: {e}]"

def read_pdf(path):
    try:
        text = ""
        with fitz.open(path) as doc:
            for page in doc:
                text += page.get_text()
        return text
    except Exception as e:
        return f"[ERROR reading PDF: {e}]"

def read_docx(path):
    try:
        doc = Document(path)
        return "\n".join([p.text for p in doc.paragraphs])
    except Exception as e:
        return f"[ERROR reading DOCX: {e}]"

def read_ipynb(path):
    try:
        with open(path, "r", encoding="utf-8") as f:
            nb = nbformat.read(f, as_version=4)
        return "\n\n".join(cell.source for cell in nb.cells if cell.cell_type in ["markdown", "code"])
    except Exception as e:
        return f"[ERROR reading IPYNB: {e}]"

def read_pptx(path):
    try:
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n\n".join(text)
    except Exception as e:
        return f"[ERROR reading PPTX: {e}]"

def read_csv(path):
    try:
        df = pd.read_csv(path, encoding='utf-8', errors='ignore')
        return df.to_string(index=False)
    except Exception as e:
        return f"[ERROR reading CSV: {e}]"

def read_excel(path):
    try:
        df = pd.read_excel(path, engine='openpyxl')
        return df.to_string(index=False)
    except Exception as e:
        return f"[ERROR reading XLSX: {e}]"

def read_md(path):
    return read_text_file(path)

def read_html(path):
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f, "html.parser")
        return soup.get_text()
    except Exception as e:
        return f"[ERROR reading HTML: {e}]"

# === Load metadata ===
with open(METADATA_PATH, "r", encoding="utf-8") as f:
    files = json.load(f)

# === Extract content ===
for file in files:
    ext = file["extension"]
    file_path = BASE_DIR / file["path"]

    if ext in {".txt", ".py", ".cpp", ".h"}:
        file["text"] = read_text_file(file_path)
    elif ext == ".pdf":
        file["text"] = read_pdf(file_path)
    elif ext == ".docx":
        file["text"] = read_docx(file_path)
    elif ext == ".ipynb":
        file["text"] = read_ipynb(file_path)
    elif ext == ".pptx":
        file["text"] = read_pptx(file_path)
    elif ext == ".csv":
        file["text"] = read_csv(file_path)
    elif ext == ".xlsx":
        file["text"] = read_excel(file_path)
    elif ext == ".md":
        file["text"] = read_md(file_path)
    elif ext == ".html":
        file["text"] = read_html(file_path)
    else:
        file["text"] = "[Unsupported file type]"

    print(f" Processed: {file['path']}")

# === Save output ===
with open(OUTPUT_PATH, "w", encoding="utf-8") as f:
    json.dump(files, f, indent=2)

print(f"\n Content extracted for {len(files)} files")
print(f" Saved to: {OUTPUT_PATH}")
